"""Comprehensive file-type coverage: docx, pptx, pdf, xlsx, png, doc, txt, csv, md.

For each type, build a file with a distinctive content token, a distinctive filename
token, and a distinctive parent-directory (path) token, then assert sparse search finds
it by content, by filename, and by path — plus that results carry the original abs path.
"""

import shutil

import pytest

from infogrep.config import Config
from infogrep.engine import SearchEngine
from infogrep.indexer import Indexer


def _pyserini_available() -> bool:
    try:
        from infogrep.jvm import ensure_jdk

        ensure_jdk()
        import pyserini.search.lucene  # noqa: F401

        return True
    except Exception:
        return False


needs_sparse = pytest.mark.skipif(not _pyserini_available(), reason="pyserini/JDK21 not available")
has_tesseract = shutil.which("tesseract") is not None


def _docx(path, text):
    import docx

    d = docx.Document()
    for line in text.split("\n"):
        d.add_paragraph(line)
    path.parent.mkdir(parents=True, exist_ok=True)
    d.save(str(path))


def _pptx(path, text):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = text
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))


def _pdf(path, text):
    import fitz

    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), text)
    path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(path))
    doc.close()


def _xlsx(path, text):
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    for i, word in enumerate(text.split(), start=1):
        ws.cell(row=1, column=i, value=word)
    path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(path))


def _png_with_text(path, text):
    """Render text into a PNG (so OCR can recover it)."""
    import fitz

    src = fitz.open()
    p = src.new_page()
    p.insert_text((72, 120), text, fontsize=40)
    pix = p.get_pixmap(dpi=200)
    path.parent.mkdir(parents=True, exist_ok=True)
    pix.save(str(path))
    src.close()


def _plain(path, text):
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(text)


@needs_sparse
def test_all_file_types_content_filename_path(tmp_path):
    # Each file: distinctive CONTENT word, distinctive FILENAME word, distinctive DIR word.
    # (filename/dir words are unique so a filename/path match is unambiguous.)
    builders = [
        ("docx", _docx, "alphacontent", "betafile", "gammadir"),
        ("pptx", _pptx, "deltacontent", "epsilonfile", "zetadir"),
        ("pdf", _pdf, "etacontent", "thetafile", "iotadir"),
        ("xlsx", _xlsx, "kappacontent", "lambdafile", "mudir"),
        ("txt", _plain, "nucontent", "xifile", "omicrondir"),
        ("csv", _plain, "picontent", "rhofile", "sigmadir"),
        ("md", _plain, "taucontent", "upsilonfile", "phidir"),
    ]
    for ext, build, content, fileword, dirword in builders:
        build(tmp_path / dirword / f"{fileword}.{ext}", content)

    cfg = Config.load(tmp_path)
    cfg.dense.enabled = False
    report = Indexer(cfg).reindex()
    assert not report.errors
    engine = SearchEngine(cfg)

    for ext, _build, content, fileword, dirword in builders:
        expected = f"{dirword}/{fileword}.{ext}"
        # by content
        hits = engine.search_sparse(content, k=5)
        assert hits and hits[0].path == expected, f"{ext}: content '{content}' -> {[h.path for h in hits]}"
        # by filename token
        hits = engine.search_sparse(fileword, k=5)
        assert any(h.path == expected for h in hits), f"{ext}: filename '{fileword}'"
        # by path/dir token
        hits = engine.search_sparse(dirword, k=5)
        assert any(h.path == expected for h in hits), f"{ext}: dir '{dirword}'"
        # result carries the original absolute path + extension
        top = engine.search_sparse(content, k=1)[0]
        assert top.abs_path == str(tmp_path / expected)
        assert top.ext == ext


@needs_sparse
def test_legacy_doc_via_textutil(tmp_path):
    if not shutil.which("textutil"):
        pytest.skip("textutil not available (non-macOS)")
    # Build a .doc by converting an .rtf with textutil.
    import subprocess

    rtf = tmp_path / "legacy" / "memo.rtf"
    rtf.parent.mkdir(parents=True)
    rtf.write_text(r"{\rtf1\ansi legacyword content about depositions.}")
    subprocess.run(["textutil", "-convert", "doc", str(rtf)], check=True)
    rtf.unlink()
    assert (tmp_path / "legacy" / "memo.doc").exists()

    cfg = Config.load(tmp_path)
    cfg.dense.enabled = False
    Indexer(cfg).reindex()
    engine = SearchEngine(cfg)
    # content extracted from .doc
    assert engine.search_sparse("legacyword depositions", k=5)[0].path == "legacy/memo.doc"
    # and findable by filename
    assert any(h.path == "legacy/memo.doc" for h in engine.search_sparse("memo", k=5))


@needs_sparse
@pytest.mark.skipif(not has_tesseract, reason="tesseract not installed")
def test_png_content_via_ocr_and_filename_without(tmp_path):
    import os

    for cand in ("/opt/homebrew/share/tessdata", "/usr/local/share/tessdata", "/usr/share/tessdata"):
        if os.path.isdir(cand):
            os.environ.setdefault("TESSDATA_PREFIX", cand)
            break

    _png_with_text(tmp_path / "figures" / "scanchart.png", "OCRWORD REVENUE CHART")

    # OCR off: png is indexed by name only -> findable by filename, not by content.
    cfg = Config.load(tmp_path)
    cfg.dense.enabled = False
    rep = Indexer(cfg).reindex()
    assert rep.name_only == 1
    engine = SearchEngine(cfg)
    assert any(h.path == "figures/scanchart.png" for h in engine.search_sparse("scanchart", k=5))
    assert engine.search_sparse("ocrword", k=5) == []  # no content without OCR

    # OCR on: rebuild -> content recoverable.
    cfg.ingest.ocr = True
    Indexer(cfg).reindex(full=True)
    hits = SearchEngine(cfg).search_sparse("ocrword revenue", k=5)
    assert hits and hits[0].path == "figures/scanchart.png"


@needs_sparse
def test_unsupported_binary_findable_by_name(tmp_path):
    (tmp_path / "data").mkdir()
    (tmp_path / "data" / "archive_zebra.bin").write_bytes(b"\x00\x01\x02\x03\xff")
    cfg = Config.load(tmp_path)
    cfg.dense.enabled = False
    cfg.include = ["**/*"]  # index the .bin too (not in the doc-only default)
    rep = Indexer(cfg).reindex()
    assert rep.name_only == 1
    # Findable by filename and by directory, even though content isn't extractable.
    engine = SearchEngine(cfg)
    assert any(h.path == "data/archive_zebra.bin" for h in engine.search_sparse("zebra", k=5))
    assert any(h.path == "data/archive_zebra.bin" for h in engine.search_sparse("data", k=5))
