"""Extractor registry: dispatch a file to the right text extractor by type.

Each extractor maps a file path to a list of :class:`ExtractedPage` and accepts (and
may ignore) the shared keyword options (``ocr``, ``ocr_min_chars``). Dispatch is by
lowercase suffix via ``_REGISTRY`` — supporting a new file type means adding one entry
there (plus its suffix to the config include defaults). Anything unregistered falls
back to a UTF-8 text reader. Files with no extractable text return an empty list — the
indexer still records them so they remain searchable by file name / path.
"""

from __future__ import annotations

import subprocess
from pathlib import Path
from typing import Callable

from ..types import ExtractedPage

Extractor = Callable[..., list[ExtractedPage]]  # (path, *, ocr=..., ocr_min_chars=...)

# Extensions treated as plain UTF-8 text (docs + common code/config formats).
_TEXT_SUFFIXES = {
    ".txt", ".md", ".markdown", ".rst", ".log", ".csv", ".tsv",
    ".json", ".jsonl", ".yaml", ".yml", ".toml", ".ini", ".cfg",
    ".html", ".htm", ".xml", ".tex",
    ".py", ".js", ".ts", ".tsx", ".jsx", ".java", ".c", ".h", ".cpp", ".hpp",
    ".cc", ".go", ".rs", ".rb", ".php", ".sh", ".bash", ".zsh", ".sql",
    ".css", ".scss", ".swift", ".kt", ".scala", ".r", ".m", ".pl",
}

# Image types: text only via OCR (requires tesseract + ingest.ocr enabled).
_IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tif", ".tiff", ".webp"}


def _extract_text(path: Path, **_opts) -> list[ExtractedPage]:
    try:
        text = path.read_text(encoding="utf-8")
    except (UnicodeDecodeError, OSError):
        return []
    return [ExtractedPage(page=None, text=text)] if text.strip() else []


def _extract_pdf(
    path: Path, ocr: bool = False, ocr_min_chars: int = 16, **_opts
) -> list[ExtractedPage]:
    import fitz  # pymupdf

    pages: list[ExtractedPage] = []
    with fitz.open(path) as doc:
        for i, page in enumerate(doc, start=1):
            text = page.get_text("text")
            # Scanned pages have little/no extractable text; OCR them if asked.
            if ocr and len(text.strip()) < ocr_min_chars:
                try:
                    tp = page.get_textpage_ocr(flags=0, full=True)
                    text = page.get_text("text", textpage=tp) or text
                except Exception:
                    pass  # tesseract missing/failed -> keep whatever we had
            if text.strip():
                pages.append(ExtractedPage(page=i, text=text))
    return pages


def _extract_docx(path: Path, **_opts) -> list[ExtractedPage]:
    import docx  # python-docx

    document = docx.Document(str(path))
    parts = [p.text for p in document.paragraphs if p.text.strip()]
    # DOCX has no fixed pagination; emit one flow page.
    text = "\n".join(parts)
    return [ExtractedPage(page=None, text=text)] if text.strip() else []


def _extract_pptx(path: Path, **_opts) -> list[ExtractedPage]:
    from pptx import Presentation  # python-pptx

    prs = Presentation(str(path))
    pages: list[ExtractedPage] = []
    for i, slide in enumerate(prs.slides, start=1):
        chunks = [
            shape.text for shape in slide.shapes if shape.has_text_frame and shape.text.strip()
        ]
        text = "\n".join(chunks)
        if text.strip():
            pages.append(ExtractedPage(page=i, text=text))
    return pages


def _extract_xlsx(path: Path, **_opts) -> list[ExtractedPage]:
    import openpyxl

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    pages: list[ExtractedPage] = []
    try:
        for i, ws in enumerate(wb.worksheets, start=1):
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    rows.append(" ".join(cells))
            text = "\n".join(rows)
            # Include the sheet title so it's searchable too.
            text = f"{ws.title}\n{text}" if text.strip() else ""
            if text.strip():
                pages.append(ExtractedPage(page=i, text=text))
    finally:
        wb.close()
    return pages


def _extract_doc(path: Path, **_opts) -> list[ExtractedPage]:
    # Legacy .doc: python-docx can't read it; use macOS `textutil` if available.
    try:
        proc = subprocess.run(
            ["textutil", "-convert", "txt", "-stdout", str(path)],
            capture_output=True, text=True, timeout=60,
        )
    except (OSError, subprocess.SubprocessError):
        return []
    text = proc.stdout or ""
    return [ExtractedPage(page=None, text=text)] if text.strip() else []


def _extract_image(path: Path, ocr: bool = False, **_opts) -> list[ExtractedPage]:
    # Images carry no text layer; only OCR yields content (needs tesseract).
    if not ocr:
        return []
    import fitz  # pymupdf

    pages: list[ExtractedPage] = []
    try:
        doc = fitz.open(path)
    except Exception:
        return []
    try:
        for i, page in enumerate(doc, start=1):
            try:
                tp = page.get_textpage_ocr(flags=0, full=True)
                text = page.get_text("text", textpage=tp)
            except Exception:
                text = ""
            if text.strip():
                pages.append(ExtractedPage(page=i, text=text))
    finally:
        doc.close()
    return pages


_REGISTRY: dict[str, Extractor] = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".doc": _extract_doc,
    ".pptx": _extract_pptx,
    ".xlsx": _extract_xlsx,
    **{suffix: _extract_image for suffix in _IMAGE_SUFFIXES},
}


def get_extractor(path: Path) -> Extractor:
    """Return the extractor for ``path`` (text reader as the default fallback)."""
    return _REGISTRY.get(path.suffix.lower(), _extract_text)


def is_supported(path: Path) -> bool:
    """Whether we have a content extractor for this file type (images need OCR)."""
    suffix = path.suffix.lower()
    return suffix in _REGISTRY or suffix in _TEXT_SUFFIXES


def extract(path: Path, ocr: bool = False, ocr_min_chars: int = 16) -> list[ExtractedPage]:
    """Extract text pages from ``path``. OCR applies to PDFs and images when enabled."""
    return get_extractor(path)(path, ocr=ocr, ocr_min_chars=ocr_min_chars)
